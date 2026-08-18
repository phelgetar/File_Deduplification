#
# Ported verbatim from doc-classifier/extract.py when that project was
# merged into File_Deduplification (2026-08-17). Unchanged so that fixes
# can still be diffed against the original.
#
#!/usr/bin/env python3
#
###################################################################
# Project: doc-classifier
# File: extract.py
# Purpose: Text extraction layer for multiple document formats
#
# Description of code and how it works:
# Converts docx, pptx, xlsx, pdf, txt, and csv files into plain text
# (with an OCR fallback for scanned PDFs) for the classifier to read
#
# Author: Tim Canady
# Created: 2026-06-17
#
# Version: 1.0.0
# Last Modified: 2026-06-17 by Tim Canady
#
# Revision History:
# - 1.0.0 (2026-06-17): Initial version
###################################################################
#
"""
Text extraction layer.

Turns a document of (almost) any common format into plain text so the
classifier brain can read it. Each format gets a small dedicated reader.
This layer is deliberately separate from the model so you can swap the
"brain" later (zero-shot -> few-shot -> RAG -> fine-tune) without touching
extraction.
"""

from __future__ import annotations

import csv
import io
from pathlib import Path

# Formats handled directly as text
PLAINTEXT_SUFFIXES = {".txt", ".md", ".markdown", ".log", ".rst", ".xml"}

# Image formats: turned into text (caption + OCR + EXIF) via vision.py
IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".heic", ".heif",
                  ".gif", ".webp", ".tiff", ".bmp"}


class ExtractionError(Exception):
    pass


def extract_text(path: str | Path) -> str:
    """Return plain text for a document, dispatching on file extension."""
    path = Path(path)
    if not path.exists():
        raise ExtractionError(f"File not found: {path}")

    suffix = path.suffix.lower()

    if suffix in PLAINTEXT_SUFFIXES:
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".csv":
        return _from_csv(path)
    if suffix == ".docx":
        return _from_docx(path)
    if suffix == ".pptx":
        return _from_pptx(path)
    if suffix == ".xlsx":
        return _from_xlsx(path)
    if suffix == ".pdf":
        return _from_pdf(path)
    if suffix in IMAGE_SUFFIXES:
        return _from_image(path)
    if suffix in {".doc", ".ppt", ".xls"}:
        raise ExtractionError(
            f"Legacy binary format '{suffix}' needs conversion. "
            "Install LibreOffice and run: "
            f"soffice --headless --convert-to {suffix[1:]}x '{path}'"
        )

    raise ExtractionError(f"Unsupported file type: {suffix}")


def _from_csv(path: Path) -> str:
    rows = []
    with path.open(newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f):
            rows.append(" ".join(row))
    return "\n".join(rows)


def _from_docx(path: Path) -> str:
    import docx  # python-docx

    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # include table cell text, which holds a lot of real content
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _from_pptx(path: Path) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def _from_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" ".join(cells))
    return "\n".join(parts)


def _from_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = [(page.extract_text() or "") for page in reader.pages]
    text = "\n".join(parts).strip()

    if not text:
        # Likely a scanned/image PDF — try OCR if optional deps are present.
        ocr = _ocr_pdf(path)
        if ocr:
            return ocr
        raise ExtractionError(
            f"No extractable text in '{path.name}'. It looks scanned. "
            "Install OCR support: pip install pymupdf pytesseract "
            "(Tesseract binary is already on this machine)."
        )
    return text


def _from_image(path: Path) -> str:
    # Imported lazily so the text-only pipeline never needs the vision deps.
    from vision import describe_image

    text = describe_image(path)
    if not text:
        raise ExtractionError(
            f"Could not describe image '{path.name}'. Is the vision model pulled? "
            "Run: ollama pull llama3.2-vision"
        )
    return text


def _ocr_pdf(path: Path) -> str:
    """Best-effort OCR. Returns '' if optional deps aren't installed."""
    try:
        import fitz  # pymupdf
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""

    parts = []
    doc = fitz.open(str(path))
    for page in doc:
        pix = page.get_pixmap(dpi=200)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        parts.append(pytesseract.image_to_string(img))
    return "\n".join(parts).strip()
